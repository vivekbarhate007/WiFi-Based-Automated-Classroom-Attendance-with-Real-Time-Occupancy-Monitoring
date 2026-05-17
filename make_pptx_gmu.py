"""
WiFi Attendance — GMU Template Presentation (revised per professor feedback)
Focus: technical challenges, design decisions, failure analysis, lessons learned.
9 slides for 8-minute presentation.
"""
from pptx import Presentation
from pptx.util import Pt

TEMPLATE = "/Users/vivekbarhate/Downloads/George Mason Presentation Templates.pptx"
OUTPUT   = "/Users/vivekbarhate/Downloads/WiFi_Attendance_GMU.pptx"

prs = Presentation(TEMPLATE)

L_TITLE   = prs.slide_layouts[0]   # Title Slide     – ph 0=title, 1=subtitle
L_CONTENT = prs.slide_layouts[1]   # Title+Content   – ph 0=title, 1=body
L_TWO_COL = prs.slide_layouts[2]   # 2-Column        – ph 0=title(right), 1=left-col, 13=right-col
L_SECTION = prs.slide_layouts[3]   # Section Header  – ph 0=title, 1=body


# ── Helpers ───────────────────────────────────────────────────────────────────

def ph(slide, idx):
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def set_title(slide, text, size=32, bold=True):
    p = ph(slide, 0)
    if not p:
        return
    tf = p.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold


def set_body(slide, ph_idx, items, base_size=18):
    """
    items: list of (text, level, bold, size_override)
      level 0 = section heading  |  level 1 = bullet
    """
    p = ph(slide, ph_idx)
    if not p:
        return
    tf = p.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for text, level, bold, sz in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.level = level
        para.space_before = Pt(4)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz if sz else base_size)
        run.font.bold = bold


def H(text, size=None, bold=True):   # heading row
    return (text, 0, bold, size)

def B(text, size=None):              # bullet row
    return (text, 1, False, size)

def GAP(size=7):                     # blank spacer
    return ("", 0, False, size)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(L_TITLE)
set_title(s1, "WiFi-Based Automated Classroom Attendance", size=30)

p = ph(s1, 1)
if p:
    tf = p.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = [
        ("Using a Commodity Android Smartphone — No Student App Required", True,  16),
        ("", False, 10),
        ("Vivek Barhate  |  Preshita Bhortake", False, 14),
        ("George Mason University  ·  April 2026", False, 13),
    ]
    first = True
    for text, bold, size in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Motivation & Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(L_CONTENT)
set_title(s2, "Motivation")
set_body(s2, 1, [
    H("Why automate attendance?", size=19),
    B("Manual roll-call in large classes wastes 5–10 min of lecture time every session"),
    B("Sign-in sheets and QR scans are trivially bypassed — a friend can sign for you"),
    B("Dedicated hardware (RFID readers, BLE beacons) costs hundreds of dollars per room"),
    GAP(),
    H("Our approach", size=19),
    B("One Android phone the professor already owns acts as the entire system"),
    B("Students use any browser — zero app installation, zero new hardware"),
    B("Attendance computed from WiFi presence: PRESENT ≥75%, PARTIAL 50–74%, ABSENT <50%"),
    GAP(),
    H("Where our effort went", size=18),
    B("Most time was spent fighting platform restrictions and designing around failure modes — not writing happy-path code"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(L_TWO_COL)
set_title(s3, "System Architecture", size=26)
set_body(s3, 1, [
    H("Student-side flow", bold=True, size=16),
    B("Connect to classroom hotspot WiFi", size=14),
    B("Open browser → captive portal at hotspot IP:8080", size=14),
    B("Enter roll number → validated against roster CSV", size=14),
    B("Success page runs JS heartbeat every 30 s", size=14),
    GAP(size=8),
    H("Professor-side flow", bold=True, size=16),
    B("Tap 'Start Session' → hotspot + NanoHTTPD server start", size=14),
    B("Ping sweep detects every connected device every 5 s", size=14),
    B("'End Session' → AttendanceEngine generates report + CSV", size=14),
], base_size=14)
set_body(s3, 13, [
    H("Components", bold=True, size=16),
    B("Router App (Android Kotlin) — hotspot + NanoHTTPD + Room DB", size=14),
    B("NanoHTTPD — embedded HTTP server, no external dependencies", size=14),
    B("FastAPI + SQLite — backend API, CSV export, web dashboard", size=14),
    B("Student browser — any browser, zero install", size=14),
    GAP(size=8),
    H("Key design choice", bold=True, size=16),
    B("Originally planned a dedicated Student App — scrapped it", size=14),
    B("Captive portal + ping sweep achieves the same without requiring students to install anything", size=13),
], base_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Technical Challenges (core slide — where effort went)
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(L_CONTENT)
set_title(s4, "Technical Challenges")
set_body(s4, 1, [
    H("Challenge 1 — ARP table blocked on Android 10+", size=18),
    B("Plan: read /proc/net/arp to find connected device IPs.  Reality: permission denied for non-system apps since Android 10"),
    B("Fix: ping-sweep all 254 IPs in the /24 subnet concurrently (semaphore=32); detect live devices by ICMP response"),
    GAP(size=6),
    H("Challenge 2 — Hotspot IP is not what we assumed", size=18),
    B("Assumed hotspot gateway = 192.168.43.1.  Physical test showed IP was 10.116.213.247"),
    B("Root cause: Android Local-Only Hotspot API assigns a random CGNAT subnet — not configurable"),
    B("Fix: read the gateway dynamically from DhcpInfo at runtime; never hardcode the IP"),
    GAP(size=6),
    H("Challenge 3 — Gson bypasses Kotlin null-safety", size=18),
    B("Kotlin data class field declared non-nullable; Gson sets it to null via reflection when JSON field is absent"),
    B("Symptom: POST /heartbeat with empty body returned HTTP 500 (expected 400)"),
    B("Fix: replaced .trim() with ?.trim() ?: \"\" — null-safe operators throughout the NanoHTTPD handlers"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Design Decisions
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(L_TWO_COL)
set_title(s5, "Design Decisions", size=26)
set_body(s5, 1, [
    H("Privacy-preserving identity", bold=True, size=16),
    B("Never store raw IPs — hash immediately: token = SHA-256(IP + sessionSalt)", size=14),
    B("sessionSalt regenerated every session → same device gets different token each class", size=14),
    B("Android 10+ blocks MAC address access for 3rd-party apps — cannot leak what we cannot read", size=14),
    GAP(size=8),
    H("Anti-proxy guards", bold=True, size=16),
    B("Problem: student A could type student B's roll number", size=14),
    B("Guard 1 — one device per check-in: token (device hash) must be unused in this session", size=14),
    B("Guard 2 — one roll number per session: roll number must not already be claimed", size=14),
    B("Both checked server-side in NanoHTTPD before any DB write", size=14),
], base_size=14)
set_body(s5, 13, [
    H("Dual-layer presence detection", bold=True, size=16),
    B("Layer 1 (IP): ping sweep every 5 s — works even if browser is closed", size=14),
    B("Layer 2 (App): JS heartbeat every 30 s from captive portal success page", size=14),
    B("Both write to presence_log; session report uses their union", size=14),
    GAP(size=8),
    H("Fault tolerance", bold=True, size=16),
    B("Brief WiFi glitches should not mark a student absent", size=14),
    B("Require 3 consecutive missed pings (~15 s) before pruning a device", size=14),
    B("Miss counter resets on the very next successful ping", size=14),
    GAP(size=8),
    H("Adaptive sampling", bold=True, size=16),
    B("Scan every 5 s during active session, 60 s when idle — saves battery", size=14),
], base_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Failure Analysis
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(L_CONTENT)
set_title(s6, "Failure Analysis — What Went Wrong & How We Fixed It")
set_body(s6, 1, [
    H("Student App planned but never shipped", size=18),
    B("Original design: dedicated Android app on each student phone sends periodic heartbeats with a stable UUID"),
    B("Problem: required every student to install it — adoption barrier too high for a classroom experiment"),
    B("Decision: captive portal browser check-in + ping sweep achieves equivalent results without installation"),
    GAP(size=6),
    H("Captive portal does not auto-open on Android", size=18),
    B("iOS and Windows auto-detect captive portals and pop up the browser — Android does not reliably do this"),
    B("Students on Android need to manually type the IP into a browser — discovered only during physical testing"),
    B("Mitigation: professor displays hotspot IP and port on screen alongside the QR code"),
    GAP(size=6),
    H("Test suite assumed wrong IP (0/8 router tests failed initially)", size=18),
    B("Automated tests hard-coded 192.168.43.1:8080 — all timed out on real device"),
    B("Logcat revealed actual IP: 10.116.213.247 — test rerun with correct host passed 7/8 immediately"),
    B("Remaining failure was the Gson null-safety bug (HTTP 500 vs expected 400) — fixed separately"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Limitations
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(L_TWO_COL)
set_title(s7, "Limitations", size=28)
set_body(s7, 1, [
    H("Platform constraints", bold=True, size=16),
    B("Android Local-Only Hotspot caps at ~10 simultaneous clients on some devices — not suitable for a 200-seat lecture hall without hardware upgrade", size=14),
    B("Hotspot SSID is randomly assigned (e.g. 'AndroidShare_4409') — professor cannot brand it or make it recognizable", size=14),
    B("JS heartbeats stop if student aggressively backgrounds the browser (especially iOS) — ping sweep acts as fallback", size=14),
    GAP(size=8),
    H("Security boundaries", bold=True, size=16),
    B("Physical proxy still possible: a student can hand their phone to a friend before leaving — no biometric or face check", size=14),
    B("Ping sweep cannot distinguish a student device from a visitor or faculty device on the same hotspot", size=14),
], base_size=14)
set_body(s7, 13, [
    H("Detection accuracy", bold=True, size=16),
    B("DHCP IP reassignment mid-session splits presence log: student's old IP and new IP are different tokens — could under-count duration", size=14),
    B("Ping RTT is a weak RSSI proxy — measures round-trip time, not signal strength; a student in the hallway with fast ping still appears present", size=14),
    GAP(size=8),
    H("Scalability & reliability", bold=True, size=16),
    B("Single point of failure: if professor's phone dies mid-session, the local Room DB session is lost (no cloud sync)", size=14),
    B("SQLite on Android is sufficient for a classroom but would need migration to a server DB for institution-wide deployment", size=14),
    B("Load test validated 35 concurrent students — real 100+ seat class would require profiling", size=14),
], base_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Evaluation & Results
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(L_CONTENT)
set_title(s8, "Evaluation & Results")
set_body(s8, 1, [
    H("Functional Tests — 20 / 20 PASS", size=18),
    B("12 FastAPI backend tests: sessions, enrollment, heartbeat, occupancy, CSV export, closed-session rejection"),
    B("8 NanoHTTPD router tests: captive portal HTML, /ping, /session, /checkin, /heartbeat, captive-portal probes"),
    B("Initial run: 7/8 router tests passed — 1 failure (500 vs 400) triggered Gson null-safety fix"),
    GAP(size=6),
    H("Load Test — 35 Concurrent Simulated Students", size=18),
    B("100% accuracy: all students enrolled, heartbeats recorded, attendance computed correctly"),
    B("p50 latency < 15 ms  ·  p95 < 50 ms  ·  p99 < 100 ms"),
    GAP(size=6),
    H("Physical Device Test (real Android phone on hotspot)", size=18),
    B("Device detected within first 5 s scan cycle  |  Ping RTT: 277 ms"),
    B("JS heartbeats confirmed in logcat on success page load"),
    B("Anti-proxy guard: second check-in from same device correctly rejected"),
    B("Departure: device pruned ~15 s after WiFi disconnect (3 missed ping cycles)"),
])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Lessons Learned
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(L_TWO_COL)
set_title(s9, "Lessons Learned", size=28)
set_body(s9, 1, [
    H("Test on real hardware early", bold=True, size=16),
    B("Every significant problem surfaced only on a physical device — not in the simulator or unit tests", size=14),
    B("The wrong IP assumption, captive portal auto-detect failure, and DHCP behaviour were all invisible until hardware testing", size=14),
    GAP(size=8),
    H("Read the platform docs, not just the tutorials", bold=True, size=16),
    B("/proc/net/arp deprecation and Local-Only Hotspot API randomization are documented — we discovered them by failing first", size=14),
    GAP(size=8),
    H("Design for failure, not just for success", bold=True, size=16),
    B("Fault tolerance (3-miss threshold), adaptive sampling, and dual-layer presence all came from asking 'what breaks next?'", size=14),
    B("Proxy guards were an afterthought — adding security late is harder and riskier than building it in", size=14),
], base_size=14)
set_body(s9, 13, [
    H("Language/framework pitfalls matter", bold=True, size=16),
    B("Gson's reflection-based deserialization silently breaks Kotlin's null-safety guarantees — a subtle bug that only appeared with a malformed request", size=14),
    B("Lesson: test error paths (bad input, missing fields) as rigorously as the happy path", size=14),
    GAP(size=8),
    H("Simplicity over complexity", bold=True, size=16),
    B("Scrapping the Student App and using a browser instead reduced the system by ~2000 lines of code and removed the biggest adoption barrier", size=14),
    B("The simpler design is more robust and more deployable", size=14),
    GAP(size=8),
    H("What we would do differently", bold=True, size=16),
    B("Add cloud sync from day one — SQLite-only means session loss if phone dies", size=14),
    B("Use a stable UUID per device (localStorage) instead of IP-hash to handle DHCP reassignment", size=14),
], base_size=14)


# ── Remove original 4 template slides ────────────────────────────────────────
xml_slides = prs.slides._sldIdLst
for _ in range(4):
    xml_slides.remove(xml_slides[0])

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            print(f"  Slide {i+1}: {shape.text_frame.text.strip()[:65]}")
            break
