from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE

# GMU colors
GMU_GREEN = RGBColor(0x00, 0x66, 0x33)
GMU_GOLD  = RGBColor(0xFF, 0xCD, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x22, 0x22, 0x22)
MID_GRAY  = RGBColor(0x55, 0x55, 0x55)

W = Inches(13.33)
H = Inches(7.5)

# Use a blank widescreen presentation — manually apply GMU colors
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

layouts = prs.slide_layouts
# Layout indices in blank template:
# 0 = Title Slide, 1 = Title and Content, 5 = Blank, 6 = Title Only


def add_textbox(slide, text, left, top, width, height,
                size=16, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullets(slide, items, left, top, width, height, size=14, color=DARK_GRAY):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "\u2022  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.3), GMU_GREEN)
    add_rect(slide, 0, Inches(1.3), W, Inches(0.06), GMU_GOLD)
    add_textbox(slide, title,
                Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.9),
                size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.9), Inches(10), Inches(0.42),
                    size=14, bold=False, color=RGBColor(0xCC, 0xFF, 0xCC))
    add_rect(slide, 0, Inches(7.2), W, Inches(0.3), GMU_GREEN)


def two_col(slide, l_title, l_items, r_title, r_items, top=Inches(1.55)):
    cw = Inches(5.9)
    gap = Inches(0.3)
    lx = Inches(0.5)
    rx = lx + cw + gap

    for x, title, items in [(lx, l_title, l_items), (rx, r_title, r_items)]:
        add_rect(slide, x, top, cw, Inches(0.4), GMU_GREEN)
        add_textbox(slide, title,
                    x + Inches(0.1), top + Inches(0.05),
                    cw - Inches(0.2), Inches(0.33),
                    size=14, bold=True, color=WHITE)
        add_bullets(slide, items,
                    x, top + Inches(0.46),
                    cw, Inches(5.2), size=14)


# ── SLIDE 1: Title ──────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(layouts[6])
bg = s1.background
bg.fill.solid()
bg.fill.fore_color.rgb = GMU_GREEN

add_rect(s1, 0, Inches(2.8), W, Inches(0.08), GMU_GOLD)
add_rect(s1, 0, Inches(6.0), W, Inches(0.08), GMU_GOLD)

add_textbox(s1, "WiFi-Based Automated\nClassroom Attendance",
            Inches(0.8), Inches(1.1), Inches(11.5), Inches(1.6),
            size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "Using a Commodity Android Smartphone — No Student App Required",
            Inches(0.8), Inches(3.05), Inches(11.5), Inches(0.55),
            size=17, color=GMU_GOLD, align=PP_ALIGN.CENTER)
add_textbox(s1, "Vivek Barhate  |  Preshita Bhortake",
            Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
            size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "George Mason University",
            Inches(0.8), Inches(4.35), Inches(11.5), Inches(0.45),
            size=16, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)
add_textbox(s1, "April 2026",
            Inches(0.8), Inches(4.88), Inches(11.5), Inches(0.4),
            size=14, color=RGBColor(0xAA, 0xCC, 0xAA), align=PP_ALIGN.CENTER)

# ── SLIDE 2: Problem ─────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(layouts[6])
header_bar(s2, "The Problem with Manual Attendance")

add_bullets(s2, [
    "Time-consuming in large lecture halls — calling roll takes 5-10 minutes per class",
    "Sign-in sheets are easily forged — proxy marking is widespread and undetectable",
    "QR-code check-in scans one device at a time — does not scale to 100+ students",
    "Dedicated hardware (RFID readers, BLE beacons) costs hundreds of dollars per room",
], Inches(0.5), Inches(1.58), Inches(12.3), Inches(3.3), size=17)

add_rect(s2, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.72), RGBColor(0xE8, 0xF5, 0xEC))
add_rect(s2, Inches(0.5), Inches(5.2), Inches(0.08), Inches(1.72), GMU_GREEN)
add_textbox(s2, "Our Goal",
            Inches(0.75), Inches(5.3), Inches(11.8), Inches(0.38),
            size=15, bold=True, color=GMU_GREEN)
add_textbox(s2, "Automate attendance using a single commodity Android phone already in the professor's pocket — zero extra hardware, zero student app, sub-30-second detection latency.",
            Inches(0.75), Inches(5.72), Inches(11.8), Inches(1.0),
            size=14, color=DARK_GRAY)

# ── SLIDE 3: Architecture ────────────────────────────────────────────────────
s3 = prs.slides.add_slide(layouts[6])
header_bar(s3, "System Architecture", "One professor phone, any student browser")

two_col(s3,
    "How It Works", [
        "Professor's phone creates a WiFi hotspot",
        "Embedded HTTP server (NanoHTTPD) serves captive portal on port 8080",
        "Student opens browser, enters roll number — check-in recorded",
        "Ping sweep every 5 s tracks all connected devices",
        "JS heartbeat from success page every 30 s — application-layer signal",
        "Session ends -> AttendanceEngine computes PRESENT / PARTIAL / ABSENT",
    ],
    "Key Components", [
        "Router App (Android) — hotspot + NanoHTTPD + Room DB",
        "FastAPI Backend (Python) — session API, CSV export, dashboard",
        "Student Browser — any browser, zero install",
        "Anti-proxy Guards — 1 device + 1 roll number per session",
        "Adaptive Sampling — 5 s active / 60 s idle scan interval",
        "Fault Tolerance — 3 missed pings (~15 s) before pruning device",
    ],
)

# ── SLIDE 4: Student Flow ─────────────────────────────────────────────────────
s4 = prs.slides.add_slide(layouts[6])
header_bar(s4, "Student Check-in Flow", "No app installation required")

steps = [
    ("1", "Connect to the classroom hotspot WiFi (SSID shown on professor's screen / QR)"),
    ("2", "Browser navigates to captive portal at http://10.x.x.x:8080"),
    ("3", "Student enters roll number — validated against pre-loaded roster CSV"),
    ("4", "Anti-proxy check: device token must be unique AND roll number must be unclaimed"),
    ("5", "Success page opens — JavaScript heartbeat POSTs to /heartbeat every 30 s"),
    ("6", "Ping sweep confirms device stays on subnet; presence_log updated every 5 s"),
]
step_top = Inches(1.55)
step_h = Inches(0.79)
for num, text in steps:
    add_rect(s4, Inches(0.4), step_top + Inches(0.08), Inches(0.45), Inches(0.45), GMU_GREEN)
    add_textbox(s4, num, Inches(0.4), step_top + Inches(0.06),
                Inches(0.45), Inches(0.45), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s4, text, Inches(1.05), step_top,
                Inches(11.8), Inches(0.72), size=15, color=DARK_GRAY)
    step_top += step_h

# ── SLIDE 5: Privacy ──────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(layouts[6])
header_bar(s5, "Privacy by Design")

two_col(s5,
    "What We Do NOT Store", [
        "IP addresses (discarded after hashing)",
        "MAC addresses (Android 10+ blocks ARP table for 3rd-party apps)",
        "Device fingerprints or browser information",
        "Cross-session linking data",
        "Location or movement data",
    ],
    "What We DO Store", [
        "SHA-256(IP + sessionSalt) — session-scoped anonymous token",
        "Roll number + name (self-provided by student)",
        "First-seen and last-seen timestamps",
        "Ping RTT as signal quality proxy (rolling 10-sample median)",
        "New random salt every session — tokens unlink across sessions",
    ],
)

add_rect(s5, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.97), RGBColor(0xFF, 0xF8, 0xDC))
add_rect(s5, Inches(0.5), Inches(6.05), Inches(0.08), Inches(0.97), GMU_GOLD)
add_textbox(s5,
    "Anti-Proxy Guards: (1) One device per check-in per session   (2) One roll number per session — both enforced server-side before any database write",
    Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.85),
    size=13, color=DARK_GRAY)

# ── SLIDE 6: Technical Deep Dive ──────────────────────────────────────────────
s6 = prs.slides.add_slide(layouts[6])
header_bar(s6, "Under the Hood")

techs = [
    ("Ping Sweep",
     "ICMP scan across entire /24 subnet every 5 s. ARP table blocked on Android 10+ for non-system apps, so we ping all 254 IPs concurrently (semaphore=32). RTT stored as rolling 10-sample median — our RSSI proxy."),
    ("Adaptive Sampling",
     "HotspotManager.setSessionActive(true) switches scan interval from 60 s (idle) to 5 s (active session). Battery overhead is negligible outside class hours."),
    ("Fault Tolerance",
     "3 consecutive missed pings (~15 s) required before a device is pruned. Handles brief WiFi glitches without false departures. Miss counter resets on the very next successful ping."),
    ("Dual-Layer Presence Detection",
     "IP-layer (ping sweep) + application-layer (JS heartbeats every 30 s from the captive portal success page). Both write to presence_log; the session report uses the union for maximum reliability."),
]
tech_top = Inches(1.58)
for title, desc in techs:
    add_rect(s6, Inches(0.4), tech_top, Inches(0.06), Inches(0.88), GMU_GREEN)
    add_textbox(s6, title,
                Inches(0.6), tech_top,
                Inches(12.2), Inches(0.36), size=15, bold=True, color=GMU_GREEN)
    add_textbox(s6, desc,
                Inches(0.6), tech_top + Inches(0.34),
                Inches(12.2), Inches(0.55), size=13, color=MID_GRAY)
    tech_top += Inches(1.24)

# ── SLIDE 7: Results ──────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(layouts[6])
header_bar(s7, "Evaluation & Results")

panels = [
    ("Functional Tests", "20 / 20 PASS", [
        "12 backend tests (FastAPI)",
        "8 router tests (NanoHTTPD)",
        "Gson null-safety fix verified (500 -> 400)",
        "Both proxy guards tested on device",
    ]),
    ("Load Test  —  35 Students", "100% Accuracy", [
        "Concurrent enroll + heartbeat rounds",
        "p50 latency  < 15 ms",
        "p95 latency  < 50 ms",
        "p99 latency  < 100 ms",
    ]),
    ("Physical Device Test", "< 5 s Detection", [
        "Detected in first scan cycle",
        "Ping RTT: 277 ms",
        "JS heartbeats confirmed in logcat",
        "Departure detected in < 15 s",
    ]),
]
panel_w = Inches(3.9)
panel_left = Inches(0.4)
for title, headline, items in panels:
    add_rect(s7, panel_left, Inches(1.58), panel_w, Inches(0.52), GMU_GREEN)
    add_textbox(s7, title,
                panel_left + Inches(0.1), Inches(1.6),
                panel_w - Inches(0.2), Inches(0.44),
                size=14, bold=True, color=WHITE)
    add_textbox(s7, headline,
                panel_left, Inches(2.22), panel_w, Inches(0.65),
                size=22, bold=True, color=GMU_GREEN, align=PP_ALIGN.CENTER)
    add_bullets(s7, items,
                panel_left + Inches(0.1), Inches(3.0),
                panel_w - Inches(0.2), Inches(3.8), size=13)
    panel_left += panel_w + Inches(0.27)

# ── SLIDE 8: Takeaways ────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(layouts[6])
header_bar(s8, "Key Takeaways & Lessons Learned")

two_col(s8,
    "What We Built", [
        "Working no-app attendance system using one commodity Android phone",
        "Privacy-preserving: SHA-256 tokens, randomized per session",
        "Dual-layer detection: ping sweep + JS heartbeats",
        "Anti-proxy guards: device + roll number uniqueness enforced",
        "Zero additional hardware cost",
        "20/20 tests passing; 35-student load test at 100% accuracy",
    ],
    "What We Learned", [
        "Android API restrictions force creative solutions — ARP blocked, so ping sweep",
        "Gson bypasses Kotlin null safety via reflection — caused 500 bug, fixed with ?.",
        "Dual-layer presence more robust than single signal alone",
        "Security must be designed in from the start — proxy guards added late",
        "Real hardware reveals gaps invisible in simulation (hotspot IP, captive portal auto-detect)",
        "Hotspot IP is random (10.x.x.x) — never assume 192.168.43.1",
    ],
)

add_rect(s8, Inches(3.2), Inches(6.1), Inches(6.9), Inches(0.82), GMU_GREEN)
add_textbox(s8, "Ready for pilot deployment in a real classroom",
            Inches(3.2), Inches(6.14), Inches(6.9), Inches(0.7),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = "/Users/vivekbarhate/Downloads/WiFi_Attendance_Presentation.pptx"
prs.save(out)
print(f"Saved to: {out}")
print(f"Total slides: {len(prs.slides)}")
